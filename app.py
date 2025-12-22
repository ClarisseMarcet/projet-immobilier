import streamlit as st
import base64
from io import BytesIO
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
import pandas as pd
from pptx import Presentation

# ==========================================================
# CONFIG STREAMLIT
# ==========================================================
st.set_page_config(
    page_title="Observatoire Immobilier & Risques Climatiques",
    layout="wide"
)

# ==========================================================
# MODE SOMBRE / MODE CLAIR



# HEADER

st.markdown("""
    <div style="
        background: linear-gradient(90deg, #0A2A43, #123B57);
        padding:15px; border-radius:8px; color:white;
        font-size:27px; font-weight:700; margin-bottom:20px;
    ">
    Observatoire Immobilier & Risques Climatiques – France
    </div>
""", unsafe_allow_html=True)

#
# EXPORTATION PDF & PPT

def generate_pdf(text):
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=A4)
    text_obj = c.beginText(40, 800)
    text_obj.setFont("Helvetica", 11)
    text_obj.setLeading(14)

    for line in text.split("\n"):
        text_obj.textLine(line[:120])

    c.drawText(text_obj)
    c.save()
    buffer.seek(0)
    return buffer


def generate_ppt(title, paragraphs):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = paragraphs[0]

    for p in paragraphs[1:]:
        para = tf.add_paragraph()
        para.text = p
        para.level = 1

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# EXPORT — SIDEBAR

st.sidebar.markdown("### Exportation")
export_text = st.sidebar.text_area("Texte à exporter (PDF / PPT)", height=180)

col1, col2 = st.sidebar.columns(2)

with col1:
    if st.button("PDF"):
        if export_text.strip():
            pdf = generate_pdf(export_text)
            b64 = base64.b64encode(pdf.read()).decode()
            st.sidebar.markdown(
                f'<a href="data:application/pdf;base64,{b64}" download="rapport.pdf">Télécharger le PDF</a>',
                unsafe_allow_html=True
            )
        else:
            st.sidebar.warning("Veuillez écrire un texte.")

with col2:
    if st.button("PPTX"):
        if export_text.strip():
            ppt = generate_ppt("Synthèse Rapport", export_text.split("\n"))
            b64 = base64.b64encode(ppt.read()).decode()
            st.sidebar.markdown(
                f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="rapport.pptx">Télécharger le PPTX</a>',
                unsafe_allow_html=True
            )
        else:
            st.sidebar.warning("Veuillez écrire un texte.")


# MESSAGE D’ACCUEIL

st.markdown("""
### Bienvenue sur l’Observatoire Immobilier & Risques Climatiques  
Sélectionnez une page dans le menu à gauche.
""")

# FOOTER
st.markdown("---")
st.markdown(
    "<center><i>Dashboard conçu par Eli  – Observatoire Immobilier & Risques Climatiques</i></center>",
    unsafe_allow_html=True
)


# PRÉSENTATION PERSONNELLE

st.markdown("<br><br><br>", unsafe_allow_html=True)
st.markdown("---")

st.markdown("""
### Présentation – Eli Clarisse AMOU

Je suis **Eli Clarisse AMOU**, étudiante en **Data Mining** à l’Université Gustave Eiffel.  
Passionnée par l’analyse de données, les statistiques et la visualisation,  
je transforme les données en outils d’aide à la décision.

Je suis actuellement **à la recherche d’un stage ou d’une alternance**  
dans les métiers de la Data (Data Analyst, Data Science, BI).

LinkedIn :  
https://www.linkedin.com/in/a-eli
""")
